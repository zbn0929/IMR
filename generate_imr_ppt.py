# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("IECE_IMR_v5_comparison.pptx").resolve()
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "ink": RGBColor(15, 23, 42),
    "muted": RGBColor(100, 116, 139),
    "line": RGBColor(226, 232, 240),
    "bg": RGBColor(248, 250, 252),
    "blue": RGBColor(37, 99, 235),
    "blue_light": RGBColor(219, 234, 254),
    "green": RGBColor(22, 163, 74),
    "green_light": RGBColor(220, 252, 231),
    "orange": RGBColor(234, 88, 12),
    "orange_light": RGBColor(255, 237, 213),
    "white": RGBColor(255, 255, 255),
}


def set_font(run, size=16, bold=False, color="ink"):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]


def add_text(slide, text, x, y, w, h, size=16, bold=False, color="ink", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = f"- {item}"
        set_font(r, size, False, color)
    return box


def add_card(slide, x, y, w, h, title, lines, accent="blue"):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLORS[accent]
    strip.line.fill.background()

    add_text(slide, title, x + 0.22, y + 0.18, w - 0.4, 0.35, 15, True, accent)
    add_bullets(slide, lines, x + 0.22, y + 0.72, w - 0.45, h - 0.86, 12.5)


def add_title(slide, title, subtitle=None, no=1):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.16))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["blue"]
    bar.line.fill.background()
    add_text(slide, title, 0.62, 0.38, 11.7, 0.55, 25, True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.0, 11.5, 0.38, 12.5, False, "muted")
    add_text(slide, f"{no:02d}", 12.1, 7.05, 0.55, 0.25, 10, False, "muted", PP_ALIGN.RIGHT)


def new_slide(no, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle, no)
    return slide


def add_flow(slide, labels, x, y, box_w=1.85, gap=0.36):
    for i, label in enumerate(labels):
        left = x + i * (box_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(y), Inches(box_w), Inches(0.68)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["blue_light"] if i < len(labels) - 1 else COLORS["green_light"]
        shape.line.color.rgb = COLORS["blue"] if i < len(labels) - 1 else COLORS["green"]
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_font(r, 11.5, True)
        if i < len(labels) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(left + box_w + 0.05), Inches(y + 0.21), Inches(0.25), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["muted"]
            arrow.line.fill.background()


def add_table(slide, data, x, y, w, h, font_size=11.5):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for col in range(cols):
        table.columns[col].width = Inches(w / cols)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = COLORS["blue"]
            else:
                cell.fill.fore_color.rgb = COLORS["white"] if r % 2 else COLORS["bg"]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_font(run, font_size, r == 0, "white" if r == 0 else "ink")


def add_note(slide, text):
    add_text(slide, f"讲稿提示：{text}", 0.66, 6.72, 11.8, 0.3, 10.5, False, "muted")


slide = new_slide(
    1,
    "IECE、IMR 与 v5 模型对比",
    "只围绕任务、原始模型和当前改进版讲清楚：为什么改、怎么改、提升了什么",
)
add_flow(slide, ["IECE 任务", "原始 IMR", "发现不足", "v5 改进", "实验结果"], 1.05, 3.25)
add_note(slide, "开场可以说：我的工作不是换一个新任务，而是在 IECE 任务上改进原始 IMR。")

slide = new_slide(2, "IECE 是什么？", "Implicit Emotion Cause Extraction：隐式情绪原因抽取")
add_card(
    slide,
    0.75,
    1.45,
    3.75,
    3.85,
    "输入",
    ["一段含有隐式情绪的文本", "文本中的候选事件", "情绪没有明显情绪词，需要推理"],
    "blue",
)
add_card(
    slide,
    4.8,
    1.45,
    3.75,
    3.85,
    "输出",
    ["判断文本整体情绪类别", "判断每个候选事件是不是原因", "核心评价通常看 IECE 的 P、R、F1"],
    "green",
)
add_card(
    slide,
    8.85,
    1.45,
    3.75,
    3.85,
    "难点",
    ["原因不一定有明显标志", "单看候选事件容易误判", "同一篇文本里多个事件会互相竞争"],
    "orange",
)
add_note(slide, "先把任务说清楚：IECE 不是普通情绪分类，而是要找出隐式情绪背后的原因事件。")

slide = new_slide(3, "原始 IMR 怎么做？", "Iterative Mutual Refinement：情绪和原因双向互相修正")
add_flow(slide, ["全文文本", "情绪分支", "情绪状态", "交叉注意力", "原因状态", "原因预测"], 0.85, 1.75)
add_flow(slide, ["候选事件", "原因分支", "原因状态", "交叉注意力", "情绪状态", "情绪预测"], 0.85, 3.05)
add_bullets(
    slide,
    [
        "IMR 的核心思想：情绪判断和原因判断不是分开的，而是互相帮助。",
        "情绪分支提供全局情绪信息，原因分支提供候选事件信息。",
        "多轮交叉注意力让两个分支不断互相修正。",
        "原始论文 IECE 结果：P=72.03，R=85.37，F1=77.94。",
    ],
    1.05,
    4.45,
    11.0,
    1.4,
)
add_note(slide, "这里要肯定 IMR：它已经解决了情绪和原因单向建模不够的问题。")

slide = new_slide(4, "原始 IMR 还存在什么不足？", "主要不足集中在原因事件建模还不够充分")
add_card(
    slide,
    0.75,
    1.45,
    3.75,
    4.15,
    "不足 1",
    ["原因侧主要依赖候选事件表示", "事件脱离全文后，语义可能不完整", "隐式情绪场景尤其需要上下文"],
    "orange",
)
add_card(
    slide,
    4.8,
    1.45,
    3.75,
    4.15,
    "不足 2",
    ["候选事件之间有篇章位置关系", "原因、背景、结果常有结构差异", "原始 IMR 没有显式加入中心事件位置先验"],
    "orange",
)
add_card(
    slide,
    8.85,
    1.45,
    3.75,
    4.15,
    "不足 3",
    ["原因预测和最终情绪匹配仍可加强", "同一事件对不同情绪的解释力不同", "需要让原因判断更受情绪状态约束"],
    "orange",
)
add_note(slide, "过渡句：所以 v5 的改进重点不是推翻 IMR，而是加强原因侧。")

slide = new_slide(5, "v5 相比 IMR 改了什么？", "保留 IMR 双向互炼主干，新增三个原因侧增强模块")
add_table(
    slide,
    [
        ["对比项", "原始 IMR", "当前 v5"],
        ["原因输入", "候选事件单视角", "event-only + text-event 双视角"],
        ["事件结构", "主要由模型隐式学习", "显式加入中心事件结构先验"],
        ["情绪约束", "交叉注意力中互相影响", "输出前再做情绪条件化原因预测"],
        ["训练策略", "常规分类损失", "Focal loss + center loss + 动态阈值"],
    ],
    0.85,
    1.45,
    11.65,
    3.2,
)
add_bullets(
    slide,
    [
        "一句话总结：v5 让模型判断原因时，同时看事件本身、事件上下文、事件位置和情绪匹配。",
        "这几个改动都围绕同一个目标：提高原因事件判断的可靠性。",
    ],
    1.05,
    5.15,
    11.0,
    0.9,
)
add_note(slide, "这一页是总览，后面三页分别讲三个核心改进。")

slide = new_slide(6, "改进一：原因侧双视角建模", "解决“只看事件本身不够”的问题")
add_card(
    slide,
    0.85,
    1.45,
    5.45,
    3.85,
    "怎么做",
    ["event-only 分支：只编码候选事件", "text-event 分支：编码全文和候选事件拼接信息", "通过 StateFusionGate 自动融合两个视角"],
    "blue",
)
add_card(
    slide,
    6.95,
    1.45,
    5.45,
    3.85,
    "解决什么",
    ["保留事件自身的局部语义", "补充事件在全文中的作用", "让模型自己决定更相信哪一种信息"],
    "green",
)
add_note(slide, "通俗讲：既看“这个事件说了什么”，也看“它放在全文里是什么意思”。")

slide = new_slide(7, "改进二：中心事件结构先验", "解决“事件位置结构没有被显式利用”的问题")
add_table(
    slide,
    [
        ["结构特征", "含义"],
        ["rel_pos", "候选事件在事件序列中的相对位置"],
        ["from_start", "距离文本开头有多远"],
        ["from_end", "距离文本结尾有多远"],
        ["log_count", "当前文本一共有多少候选事件"],
    ],
    0.9,
    1.45,
    5.5,
    2.35,
)
add_card(
    slide,
    6.85,
    1.45,
    5.55,
    2.35,
    "怎么做",
    ["把 4 个结构特征编码成 prior_state", "轻微修正 cause_state", "同时输出 center_logits，作为辅助监督信号"],
    "blue",
)
add_bullets(
    slide,
    [
        "它不是硬规则，不会强行规定某个位置一定是原因。",
        "它只是给模型一个弱提示：哪些候选事件更像文本的中心事件。",
        "训练中增加 center loss，权重为 0.2；最终 cause logits 加入 0.15 * center_logits。",
    ],
    1.05,
    4.35,
    11.2,
    1.25,
)
add_note(slide, "讲的时候强调“结构先验是辅助信息，不是直接改标签”。")

slide = new_slide(8, "改进三：情绪条件化原因预测", "解决“原因判断和最终情绪匹配不够显式”的问题")
add_flow(slide, ["情绪状态", "条件化模块", "修正原因状态", "原因分类器", "原因预测"], 1.4, 2.0)
add_bullets(
    slide,
    [
        "原始 IMR 通过交叉注意力让情绪和原因互相影响。",
        "v5 在输出原因预测前，再把最终情绪状态融合进原因状态。",
        "直观理解：判断一个事件是不是原因时，再问一次“它和当前情绪搭不搭”。",
    ],
    1.05,
    3.65,
    11.0,
    1.45,
)
add_note(slide, "这页不用讲复杂公式，重点讲“原因必须能解释情绪”。")

slide = new_slide(9, "结果对比", "当前可确认的 IECE 10 折平均 F1")
add_table(
    slide,
    [
        ["模型", "主要特点", "IECE F1", "相对 IMR"],
        ["原始 IMR", "情绪-原因双向互炼", "77.94", "-"],
        ["v5", "双视角原因建模 + 中心事件先验 + 情绪条件化", "79.19", "+1.25"],
    ],
    1.05,
    1.55,
    11.2,
    1.9,
    12.5,
)
add_card(
    slide,
    1.05,
    4.05,
    11.2,
    1.35,
    "结论",
    [
        "v5 在原始 IMR 基础上提升约 1.25 个 F1 点。",
        "提升来源主要是原因侧信息更充分：事件本身、上下文、结构位置和情绪匹配共同参与判断。",
    ],
    "green",
)
add_note(slide, "如果老师问阈值：v5 最优结果使用验证集动态阈值，固定 0.5 是后续公平性补充实验。")

slide = new_slide(10, "汇报总结", "三句话讲清楚这版工作的价值")
add_card(
    slide,
    0.9,
    1.4,
    11.5,
    1.2,
    "1. IECE 的难点",
    ["隐式情绪没有明显情绪词，原因事件需要结合全文和篇章结构推理。"],
    "blue",
)
add_card(
    slide,
    0.9,
    3.0,
    11.5,
    1.2,
    "2. IMR 的基础",
    ["原始 IMR 通过情绪分支和原因分支的双向互炼，解决了情绪和原因互相依赖的问题。"],
    "green",
)
add_card(
    slide,
    0.9,
    4.6,
    11.5,
    1.35,
    "3. v5 的提升",
    ["v5 重点增强原因侧，让模型同时利用事件本身、上下文、中心事件结构先验和情绪匹配，最终 F1 从 77.94 提升到 79.19。"],
    "orange",
)
add_note(slide, "最后落点：v5 是对 IMR 的原因侧增强，不是换任务，也不是只调参。")


prs.save(OUT)
print(OUT)
